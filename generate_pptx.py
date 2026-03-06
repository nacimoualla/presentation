from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# Dracula Theme Colors
BG_COLOR = RGBColor(40, 42, 54)
TEXT_COLOR = RGBColor(248, 248, 242)
HIGHLIGHT_COLOR = RGBColor(255, 121, 198)
SUCCESS_COLOR = RGBColor(80, 250, 123)
DANGER_COLOR = RGBColor(255, 85, 85)

def style_text_frame(text_frame, font_size=None, color=TEXT_COLOR, bold=False, align=None):
    for paragraph in text_frame.paragraphs:
        if align is not None:
            paragraph.alignment = align
        for run in paragraph.runs:
            run.font.color.rgb = color
            run.font.name = 'Segoe UI'
            if font_size:
                run.font.size = Pt(font_size)
            run.font.bold = bold

def style_shape(shape, font_size=None, color=TEXT_COLOR, bold=False, align=None):
    if not shape.has_text_frame:
        return
    style_text_frame(shape.text_frame, font_size, color, bold, align)

def set_slide_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

# ----------------- 1. Title Slide -----------------
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
set_slide_bg(slide)

title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Multi-Threading & Synchronisation"
subtitle.text = "Maîtriser l'exécution parallèle et éviter le chaos.\n\nPrésenté par : Nacim Oualla, Essaqy Douae, Yassmine Mouden, Ouachikh Zakaria"

style_shape(title, font_size=44, color=HIGHLIGHT_COLOR, bold=True)
style_shape(subtitle, font_size=18, color=TEXT_COLOR)

# ----------------- Helper for Bullet Slides -----------------
bullet_slide_layout = prs.slide_layouts[1]

def add_styled_slide(title_text, bullet_points):
    slide = prs.slides.add_slide(bullet_slide_layout)
    set_slide_bg(slide)
    
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = title_text
    style_shape(title_shape, font_size=36, color=HIGHLIGHT_COLOR, bold=True)
    
    tf = body_shape.text_frame
    tf.clear() 
    
    for pt in bullet_points:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = pt
        run.font.color.rgb = TEXT_COLOR
        run.font.name = 'Segoe UI'
        run.font.size = Pt(22)
        p.space_after = Pt(14)
        
    return slide

# ----------------- 2. What is a Thread -----------------
add_styled_slide("Qu'est-ce qu'un Thread ?", [
    "Un processus est une application en cours d'exécution.",
    "Un thread est la plus petite unité de traitement.",
    "Attention : Les threads partagent le même espace mémoire !"
])

# ----------------- 3. Processus vs Thread table -----------------
slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only layout
set_slide_bg(slide)
slide.shapes.title.text = "Processus vs Thread"
style_shape(slide.shapes.title, font_size=36, color=HIGHLIGHT_COLOR, bold=True)

rows = 5
cols = 2
left = Inches(0.5)
top = Inches(1.8)
width = Inches(9.0)
height = Inches(3.0)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

table.cell(0, 0).text = "Processus"
table.cell(0, 1).text = "Thread"

data = [
    ("Espaces mémoire isolés", "Mémoire partagée"),
    ("Lourd (Temps CPU)", "Léger et rapide"),
    ("Forte isolation", "Faible isolation (un crash tue le processus)"),
    ("Communication (IPC) lente", "Communication très rapide")
]

for row_idx, row_data in enumerate(data):
    for col_idx, text in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = text
        style_text_frame(cell.text_frame, font_size=18, color=TEXT_COLOR)

for col_idx in range(2):
    cell = table.cell(0, col_idx)
    style_text_frame(cell.text_frame, font_size=20, color=HIGHLIGHT_COLOR, bold=True)

# ----------------- 4. Pourquoi le Multi-Threading ? -----------------
add_styled_slide("Pourquoi le Multi-Threading ?", [
    "1. Réactivité : L'interface utilisateur ne gèle pas.",
    "2. Performance : Exploitation maximale des processeurs multi-cœurs.",
    "3. Économie : Créer un thread coûte beaucoup moins de mémoire."
])

# ----------------- 5. Cas d'usage Pratiques -----------------
add_styled_slide("Cas d'usage Pratiques", [
    "Jeux Vidéo : Rendu graphique, physique, IA et son (threads distincts).",
    "Serveurs Web : Apache/Tomcat gère chaque requête via un thread.",
    "Big Data : Traitement massivement parallèle, encodage vidéo, rendu 3D."
])

# ----------------- 6. Le Danger : Race Conditions -----------------
add_styled_slide("Le Danger : Race Conditions", [
    "Que se passe-t-il si deux threads modifient la même donnée ?",
    "Données corrompues",
    "Comportements imprévisibles",
    "Crash de l'application"
])

# ----------------- 7. La Synchronisation -----------------
add_styled_slide("La Synchronisation", [
    "Mettre de l'ordre dans l'exécution avec des verrous.",
    "Mutex : Verrou d'exclusion mutuelle. Un seul thread.",
    "Sémaphore : Limitation d'accès. N threads maximum.",
    "Moniteur : Approche Objet. Méthodes synchronisées nativement."
])

# ----------------- 8. Le Cauchemar : Deadlocks -----------------
add_styled_slide("Le Cauchemar : Deadlocks", [
    "L'Interblocage fatal.",
    "\"Le Thread A attend que le Thread B lâche une ressource...",
    "...mais le Thread B attend que le Thread A lâche la sienne.\"",
    "Résultat : Tout se fige."
])

# ----------------- 9. Langages & Concurrence -----------------
add_styled_slide("Langages & Concurrence", [
    "Python : Le GIL (Global Interpreter Lock) limite le parallélisme CPU natif.",
    "Java / C++ : Threads OS natifs. Performants mais complexes.",
    "Go : Goroutines ultra-légers.",
    "Node.js : Asynchronisme mono-thread via l'Event Loop."
])

# ----------------- 10. Bonnes Pratiques -----------------
add_styled_slide("Bonnes Pratiques", [
    "Minimiser les données partagées entre les threads.",
    "Garder les sections critiques (verrouillées) très courtes.",
    "Toujours utiliser unlock() dans un bloc finally."
])

# ----------------- 11. Closing slide -----------------
slide = prs.slides.add_slide(title_slide_layout)
set_slide_bg(slide)
slide.shapes.title.text = "Merci !"
slide.placeholders[1].text = "Avez-vous des questions ?"
style_shape(slide.shapes.title, font_size=50, color=SUCCESS_COLOR, bold=True)
style_shape(slide.placeholders[1], font_size=32, color=TEXT_COLOR)


prs.save("Presentation_MultiThreading_Groupe.pptx")
print("Presentation Presentation_MultiThreading_Groupe.pptx generated successfully!")
